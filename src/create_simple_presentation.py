#!/usr/bin/env python3
"""
Simple PowerPoint Generator Example
Demonstrates how to create PowerPoint presentations with python-pptx
including adding images and customizing layouts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def create_simple_presentation():
    """Create a simple presentation example."""
    prs = Presentation()
    
    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "My Test Results"
    subtitle.text = "Generated with Python"
    
    # Slide 2: Content slide with bullet points
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Test Summary"
    content.text = """• All tests passed
• Coverage: 100%
• No errors found
• Ready for deployment"""
    
    # Slide 3: Adding an image (if file exists)
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Architecture Diagram"
    
    # Try to add an image - replace with your actual image path
    diagram_path = "/home/snowholt/coding/python/saba_softwareTesting/src/docs/diagrams"
    
    # Check if any diagram images exist
    if os.path.exists(diagram_path):
        png_files = [f for f in os.listdir(diagram_path) if f.endswith('.png')]
        if png_files:
            img_path = os.path.join(diagram_path, png_files[0])
            try:
                slide.shapes.add_picture(img_path, Inches(1), Inches(2), Inches(8), Inches(4))
            except Exception as e:
                # If image fails, add text instead
                textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                textbox.text = f"Image placeholder\n(Error loading {png_files[0]}: {str(e)})"
        else:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            textbox.text = "No PNG diagrams found in diagrams folder"
    else:
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        textbox.text = "Diagrams folder not found"
    
    # Slide 4: Table example (like test results)
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Test Results Table"
    
    # Add table
    rows, cols = 4, 3
    table = slide.shapes.add_table(rows, cols, Inches(2), Inches(2), Inches(6), Inches(3)).table
    
    # Header row
    table.cell(0, 0).text = "Test Type"
    table.cell(0, 1).text = "Count"
    table.cell(0, 2).text = "Status"
    
    # Data rows
    table.cell(1, 0).text = "Unit Tests"
    table.cell(1, 1).text = "13"
    table.cell(1, 2).text = "PASS"
    
    table.cell(2, 0).text = "Integration Tests"
    table.cell(2, 1).text = "10"
    table.cell(2, 2).text = "PASS"
    
    table.cell(3, 0).text = "System Tests"
    table.cell(3, 1).text = "8"
    table.cell(3, 2).text = "PASS"
    
    # Save the presentation
    output_file = "/home/snowholt/coding/python/saba_softwareTesting/src/docs/Simple_Example.pptx"
    prs.save(output_file)
    print(f"Simple presentation saved as: {output_file}")
    return output_file


def create_presentation_with_mermaid_diagrams():
    """Create presentation that references the existing Mermaid diagrams."""
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Software Testing Project"
    subtitle.text = "With Mermaid Diagrams Integration"
    
    # Read and reference the existing diagrams
    diagrams_info = [
        ("Project Architecture", "docs/diagrams/project_architecture.md"),
        ("Application Flowchart", "docs/diagrams/application_flowchart.md"),
        ("Testing Workflow", "docs/diagrams/testing_workflow.md")
    ]
    
    for diagram_title, diagram_path in diagrams_info:
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = diagram_title
        
        # Try to read the diagram file and show some content
        full_path = f"/home/snowholt/coding/python/saba_softwareTesting/src/{diagram_path}"
        try:
            if os.path.exists(full_path):
                with open(full_path, 'r') as f:
                    diagram_content = f.read()[:500]  # First 500 chars
                content.text = f"Diagram available at: {diagram_path}\n\nPreview:\n{diagram_content}..."
            else:
                content.text = f"Diagram file not found: {diagram_path}"
        except Exception as e:
            content.text = f"Error reading diagram: {str(e)}"
    
    # Save
    output_file = "/home/snowholt/coding/python/saba_softwareTesting/src/docs/Mermaid_Diagrams_Presentation.pptx"
    prs.save(output_file)
    print(f"Mermaid diagrams presentation saved as: {output_file}")
    return output_file


if __name__ == "__main__":
    print("Creating presentations...")
    
    # Create simple example
    simple_file = create_simple_presentation()
    
    # Create mermaid diagrams presentation
    mermaid_file = create_presentation_with_mermaid_diagrams()
    
    print("\nAll presentations created successfully!")
    print(f"Files created:")
    print(f"1. {simple_file}")
    print(f"2. {mermaid_file}")
