#!/usr/bin/env python3
"""
PowerPoint Presentation Generator
Converts VBA PowerPoint automation code to Python using python-pptx library.

This script creates a comprehensive PowerPoint presentation for the Personal Finance
Calculator testing project, equivalent to the provided VBA code.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime


class PowerPointGenerator:
    """Class to generate PowerPoint presentation from testing data."""
    
    def __init__(self):
        """Initialize the PowerPoint presentation."""
        self.prs = Presentation()
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        
    def create_full_presentation(self):
        """Create the complete presentation with all slides."""
        print("Creating PowerPoint presentation...")
        
        # Add all slides in order
        self.add_title_slide()
        self.add_project_overview_slide()
        self.add_testing_strategy_slide()
        self.add_architecture_diagram_slide()
        self.add_application_workflow_slide()
        self.add_testing_workflow_slide()
        self.add_unit_testing_results_slide()
        self.add_unit_test_details_slide()
        self.add_integration_testing_results_slide()
        self.add_integration_test_cases_slide()
        self.add_system_testing_results_slide()
        self.add_system_test_cases_slide()
        self.add_overall_test_summary_slide()
        self.add_pytest_results_slide()
        self.add_test_coverage_analysis_slide()
        self.add_test_documentation_slide()
        self.add_application_features_slide()
        self.add_input_validation_system_slide()
        self.add_testing_best_practices_slide()
        self.add_error_handling_edge_cases_slide()
        self.add_performance_testing_slide()
        self.add_project_structure_slide()
        self.add_development_tools_setup_slide()
        self.add_quality_assurance_summary_slide()
        self.add_conclusion_next_steps_slide()
        
        return self.prs
    
    def add_title_slide(self):
        """Add title slide."""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Personal Finance Calculator"
        subtitle.text = "Comprehensive Software Testing Project"
        
        # Add footer with date
        footer = slide.shapes.add_textbox(
            Inches(1), Inches(8), Inches(8), Inches(0.5)
        )
        footer.text = f"Testing Documentation - {datetime.now().strftime('%B %Y')}"
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    def add_project_overview_slide(self):
        """Add project overview slide."""
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Project Overview"
        content.text = """Purpose: Demonstrate comprehensive software testing methodologies
Application: Personal Finance Calculator
Architecture: Modular Python application
Testing Approach: Three-tier testing strategy

Key Features:
• Simple Interest Calculator
• Compound Interest Calculator
• Loan Payment Calculator
• Input Validation System
• Error Handling"""
        
    def add_testing_strategy_slide(self):
        """Add testing strategy slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Testing Strategy Overview"
        content.text = """Testing Levels:
1. Unit Testing: Individual function validation
2. Integration Testing: Component interaction
3. System Testing: End-to-end workflow

Framework:
• Primary: Python unittest
• Secondary: pytest with HTML reporting

Metrics:
• Coverage: 100% function coverage
• Success Rate: 100%"""
        
    def add_architecture_diagram_slide(self):
        """Add architecture diagram slide."""
        slide_layout = self.prs.slide_layouts[5]  # Title only
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Project Architecture Diagram"
        
        # Add placeholder text box for diagram
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        textbox.text = """[Architecture Diagram Placeholder]

Main Components:
• Calculator Module (calculator.py)
• Validator Module (validator.py)
• Main Application (main.py)
• Test Suite (tests/)

Note: Diagram images can be added using slide.shapes.add_picture()"""
        
    def add_application_workflow_slide(self):
        """Add application workflow slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Application Workflow"
        
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        textbox.text = """[Application Workflow Diagram Placeholder]

Workflow Steps:
1. User Input → Validation
2. Validated Input → Calculator
3. Calculator → Results
4. Results → User Output

Note: Workflow diagrams can be added using slide.shapes.add_picture()"""
        
    def add_testing_workflow_slide(self):
        """Add testing workflow slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Testing Workflow"
        
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        textbox.text = """[Testing Workflow Diagram Placeholder]

Testing Process:
1. Unit Tests → Individual Functions
2. Integration Tests → Component Interaction
3. System Tests → End-to-End
4. Report Generation → Documentation

Note: Testing workflow diagrams available in docs/diagrams/"""
        
    def add_unit_testing_results_slide(self):
        """Add unit testing results slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Unit Testing Results"
        
        # Create a simple table
        rows, cols = 2, 4
        table = slide.shapes.add_table(
            rows, cols, Inches(2), Inches(2), Inches(6), Inches(1.5)
        ).table
        
        # Header row
        table.cell(0, 0).text = "Tests Run"
        table.cell(0, 1).text = "Failures"
        table.cell(0, 2).text = "Errors"
        table.cell(0, 3).text = "Success Rate"
        
        # Data row
        table.cell(1, 0).text = "13"
        table.cell(1, 1).text = "0"
        table.cell(1, 2).text = "0"
        table.cell(1, 3).text = "100.0%"
        
    def add_unit_test_details_slide(self):
        """Add unit test details slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Unit Test Details"
        
        # Add code box with monospace font
        code_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        code_text = """def test_simple_interest_calculation(self):
    result = self.calculator.calculate_simple_interest(1000, 5, 2)
    self.assertEqual(result, 100.0)

def test_compound_interest_calculation(self):
    result = self.calculator.calculate_compound_interest(1000, 5, 2)
    self.assertAlmostEqual(result, 102.5, places=2)"""
        
        code_box.text = code_text
        # Set font to monospace
        for paragraph in code_box.text_frame.paragraphs:
            paragraph.font.name = 'Courier New'
            paragraph.font.size = Pt(12)
            
    def add_integration_testing_results_slide(self):
        """Add integration testing results slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Integration Testing Results"
        
        rows, cols = 2, 4
        table = slide.shapes.add_table(
            rows, cols, Inches(2), Inches(2), Inches(6), Inches(1.5)
        ).table
        
        # Header row
        table.cell(0, 0).text = "Tests Run"
        table.cell(0, 1).text = "Failures"
        table.cell(0, 2).text = "Errors"
        table.cell(0, 3).text = "Success Rate"
        
        # Data row
        table.cell(1, 0).text = "10"
        table.cell(1, 1).text = "0"
        table.cell(1, 2).text = "0"
        table.cell(1, 3).text = "100.0%"
        
    def add_integration_test_cases_slide(self):
        """Add integration test cases slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Integration Test Cases"
        
        code_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        code_text = """def test_calculator_with_validation(self):
    result = self.app.process_input(1000, 5, 2, 'simple')
    self.assertEqual(result, 100.0)

def test_validator_integration(self):
    # Test validator with calculator
    valid_input = self.validator.validate_positive_number(1000)
    result = self.calculator.calculate_simple_interest(valid_input, 5, 2)
    self.assertEqual(result, 100.0)"""
        
        code_box.text = code_text
        for paragraph in code_box.text_frame.paragraphs:
            paragraph.font.name = 'Courier New'
            paragraph.font.size = Pt(12)
            
    def add_system_testing_results_slide(self):
        """Add system testing results slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "System Testing Results"
        
        rows, cols = 2, 4
        table = slide.shapes.add_table(
            rows, cols, Inches(2), Inches(2), Inches(6), Inches(1.5)
        ).table
        
        # Header row
        table.cell(0, 0).text = "Tests Run"
        table.cell(0, 1).text = "Failures"
        table.cell(0, 2).text = "Errors"
        table.cell(0, 3).text = "Success Rate"
        
        # Data row
        table.cell(1, 0).text = "8"
        table.cell(1, 1).text = "0"
        table.cell(1, 2).text = "0"
        table.cell(1, 3).text = "100.0%"
        
    def add_system_test_cases_slide(self):
        """Add system test cases slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "System Test Cases"
        
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(3)
        )
        text_box.text = """Test Case: End-to-end calculation with invalid input
Expected: Raises ValueError
Result: Passed

Test Case: Complete workflow with valid inputs
Expected: Correct calculation result
Result: Passed

Test Case: Error handling throughout system
Expected: Proper error messages
Result: Passed"""
        
    def add_overall_test_summary_slide(self):
        """Add overall test summary slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Overall Test Summary"
        content.text = """Total Tests: 31
Unit Tests: 13
Integration Tests: 10
System Tests: 8
Failures: 0
Success Rate: 100%

All test categories passed successfully!"""
        
    def add_pytest_results_slide(self):
        """Add pytest results slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Pytest Results"
        
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        text_box.text = """collected 31 items

test_unit.py ............. [41%]
test_integration.py .......... [74%]
test_system.py ........ [100%]

====== 31 passed in 0.42s ======"""
        
        for paragraph in text_box.text_frame.paragraphs:
            paragraph.font.name = 'Courier New'
            paragraph.font.size = Pt(14)
            
    def add_test_coverage_analysis_slide(self):
        """Add test coverage analysis slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Test Coverage Analysis"
        content.text = """Coverage: 100% of functions
Lines Covered: 85/85
Branches Covered: 100%
Tool: Coverage.py

Complete coverage achieved across all modules!"""
        
    def add_test_documentation_slide(self):
        """Add test documentation slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Test Documentation"
        content.text = """Documents:
• Test Plan
• Test Cases
• Test Results
• Generated: pytest-html report

All documentation available in docs/test_reports/"""
        
    def add_application_features_slide(self):
        """Add application features slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Application Features"
        content.text = """• Simple Interest Calculation
• Compound Interest Calculation
• Loan Payment Calculation
• Input Validation
• Error Handling

Comprehensive financial calculation suite with robust validation."""
        
    def add_input_validation_system_slide(self):
        """Add input validation system slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Input Validation System"
        
        code_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        code_text = """def validate_input(self, value):
    if not isinstance(value, (int, float)):
        raise ValueError('Input must be numeric')
    if value < 0:
        raise ValueError('Input must be non-negative')
    return True

def validate_positive_number(self, value):
    self.validate_input(value)
    if value <= 0:
        raise ValueError('Input must be positive')
    return value"""
        
        code_box.text = code_text
        for paragraph in code_box.text_frame.paragraphs:
            paragraph.font.name = 'Courier New'
            paragraph.font.size = Pt(11)
            
    def add_testing_best_practices_slide(self):
        """Add testing best practices slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Testing Best Practices"
        content.text = """• Write clear, concise tests
• Test edge cases
• Automate testing
• Maintain high coverage
• Document test cases
• Use meaningful test names
• Keep tests independent
• Test both positive and negative scenarios"""
        
    def add_error_handling_edge_cases_slide(self):
        """Add error handling edge cases slide."""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Error Handling & Edge Cases"
        
        code_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        code_text = """def test_negative_input(self):
    with self.assertRaises(ValueError):
        self.calculator.calculate_simple_interest(-1000, 5, 2)

def test_zero_time_period(self):
    with self.assertRaises(ValueError):
        self.calculator.calculate_simple_interest(1000, 5, 0)

def test_invalid_string_input(self):
    with self.assertRaises(ValueError):
        self.calculator.calculate_simple_interest("invalid", 5, 2)"""
        
        code_box.text = code_text
        for paragraph in code_box.text_frame.paragraphs:
            paragraph.font.name = 'Courier New'
            paragraph.font.size = Pt(11)
            
    def add_performance_testing_slide(self):
        """Add performance testing slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Performance Testing"
        content.text = """Execution Time: < 0.01s per calculation
Memory Usage: Minimal
Tested With: 10,000 iterations

Performance benchmarks meet requirements for
real-time financial calculations."""
        
    def add_project_structure_slide(self):
        """Add project structure slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Project Structure & Organization"
        content.text = """src/
• calculator.py
• validator.py
• main.py

tests/
• test_unit.py
• test_integration.py
• test_system.py

docs/
• test_reports/
• diagrams/"""
        
    def add_development_tools_setup_slide(self):
        """Add development tools setup slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Development Tools & Setup"
        content.text = """Language: Python 3.12
IDE: VS Code

Tools:
• unittest
• pytest
• coverage.py
• pytest-html
• python-pptx (for this presentation)"""
        
    def add_quality_assurance_summary_slide(self):
        """Add quality assurance summary slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Quality Assurance Summary"
        content.text = """Quality Goals Achieved:
• No defects
• High reliability
• Full test coverage
• Comprehensive documentation
• Automated testing pipeline

Project demonstrates industry-standard QA practices."""
        
    def add_conclusion_next_steps_slide(self):
        """Add conclusion and next steps slide."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Conclusion & Next Steps"
        content.text = """Achievements:
• Comprehensive testing strategy implemented
• 100% success rate across all test levels
• Complete documentation and reporting

Next Steps:
• GUI development
• Additional calculation types
• Performance optimization
• Continuous integration setup

Key Takeaway: Systematic testing ensures software quality"""
        
    def save_presentation(self, filename="PersonalFinanceCalculator_TestingPresentation.pptx"):
        """Save the presentation to file."""
        self.prs.save(filename)
        print(f"Presentation saved as: {filename}")
        return filename


def main():
    """Main function to create and save the presentation."""
    generator = PowerPointGenerator()
    presentation = generator.create_full_presentation()
    
    # Save in the docs directory
    output_path = "/home/snowholt/coding/python/saba_softwareTesting/src/docs/PersonalFinanceCalculator_TestingPresentation.pptx"
    generator.save_presentation(output_path)
    
    print(f"\nPresentation created successfully!")
    print(f"Total slides: {len(presentation.slides)}")
    print(f"Saved to: {output_path}")
    
    return output_path


if __name__ == "__main__":
    main()
